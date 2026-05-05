"""Generate filled RCA PowerPoint from template."""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from models import RCAIncident

# Black color for text
BLACK = RGBColor(0, 0, 0)


def format_datetime(dt) -> str:
    """Format datetime for display."""
    if dt is None:
        return ""
    if hasattr(dt, 'strftime'):
        return dt.strftime("%m/%d/%Y %H:%M")
    return str(dt)


def format_date(d) -> str:
    """Format date for display."""
    if d is None:
        return ""
    if hasattr(d, 'strftime'):
        return d.strftime("%m/%d/%Y")
    return str(d)


def set_cell_text(table, row: int, col: int, text: str, font_size: int = None):
    """Safely set text in a table cell with black font color.
    
    Args:
        table: The table object
        row: Row index
        col: Column index  
        text: Text to set
        font_size: Optional font size in points (for fitting text in small boxes)
    """
    try:
        cell = table.cell(row, col)
        cell.text = str(text) if text else ""
        # Set font color to black for all paragraphs/runs
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = BLACK
                if font_size:
                    run.font.size = Pt(font_size)
    except IndexError:
        pass  # Cell doesn't exist, skip


def fill_slide_1(slide, data: RCAIncident):
    """Fill Slide 1: Incident Overview."""
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        
        table = shape.table
        rows = len(table.rows)
        cols = len(table.columns)
        
        # Table 7 (1x2): Location
        # Table 6 (1x2): Job Description  
        # Table 9 (1x2): Incident Type
        if rows == 1 and cols == 2:
            first_cell = table.cell(0, 0).text.strip()
            if "Location" in first_cell and "Incident" not in first_cell:
                set_cell_text(table, 0, 1, data.site_location)
            elif "Job Description" in first_cell:
                set_cell_text(table, 0, 1, data.job_description, font_size=9)
            elif "Incident Type" in first_cell and rows == 1:
                set_cell_text(table, 0, 1, data.kind_of_injury)
        
        # Table 5 (8x2): Incident Type checkboxes
        # Row 0 is header, rows 1-7 are checkboxes
        # Labels are in column 1, X goes in column 0
        elif rows == 8 and cols == 2:
            incident_type_map = {
                "OSHA Recordable": "OSHA Recordable / Telemed",
                "Lost Time Injury": "Lost Time Injury",
                "Trailer Pullout": "Trailer Pullout",
                "Non-Medical or Medical Only": "Med Only / Nurse Triage",
                "PIT on PIT": "PIT on PIT /PIT on Structure",
                "PIT on Pedestrian": "PIT on Pedestrian / Near Miss",
            }
            for row_idx in range(1, 8):
                label = table.cell(row_idx, 1).text.strip()
                # Check if any of our incident types match this row
                matched = False
                if data.incident_types:
                    for itype in data.incident_types:
                        if itype in label or label in itype or any(k in label for k in incident_type_map if incident_type_map[k] == itype):
                            matched = True
                            break
                set_cell_text(table, row_idx, 0, "X" if matched else "")
        
        # Table 8 (4x2): Dates & Personnel
        elif rows == 4 and cols == 2:
            first_cell = table.cell(0, 0).text.strip()
            if "Date/Time" in first_cell:
                set_cell_text(table, 0, 1, format_datetime(data.incident_datetime))
                set_cell_text(table, 1, 1, format_datetime(data.reported_datetime))
                set_cell_text(table, 2, 1, format_datetime(data.medical_datetime) if data.medical_datetime else "N/A")
                set_cell_text(table, 3, 1, data.associates_manager)
        
        # Table 2 (3x2): Work Status Y/N
        elif rows == 3 and cols == 2:
            first_cell = table.cell(0, 0).text.strip()
            if "Normal Duties" in first_cell:
                set_cell_text(table, 0, 1, data.doing_normal_duties)
                set_cell_text(table, 1, 1, data.on_normal_shift)
                set_cell_text(table, 2, 1, data.on_overtime)
        
        # Table 3 (2x1): Incident Location
        elif rows == 2 and cols == 1:
            first_cell = table.cell(0, 0).text.strip()
            if "Incident Location" in first_cell:
                set_cell_text(table, 1, 0, data.incident_location)


def fill_slide_2(slide, data: RCAIncident):
    """Fill Slide 2: Associate & Incident Details."""
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        
        table = shape.table
        rows = len(table.rows)
        cols = len(table.columns)
        
        # Table 2 (12x2): Associate Work History
        if rows == 12 and cols == 2:
            set_cell_text(table, 0, 1, data.hours_in_path)
            set_cell_text(table, 1, 1, data.hours_this_week)
            set_cell_text(table, 2, 1, data.glide_path)
            set_cell_text(table, 3, 1, data.no_training_paperwork)
            set_cell_text(table, 4, 1, data.labor_share)
            set_cell_text(table, 5, 1, data.yes_training_paperwork)
            set_cell_text(table, 6, 1, data.previous_stand_down)
            set_cell_text(table, 7, 1, data.roster_reviewed)
            set_cell_text(table, 8, 1, data.countermeasure_updated)
            set_cell_text(table, 9, 1, data.last_safety_observation)
            set_cell_text(table, 10, 1, data.observation_findings)
            set_cell_text(table, 11, 1, data.disciplinary_actions)
        
        # Table 8 (4x2): Incident Category checkboxes
        elif rows == 4 and cols == 2:
            for row_idx in range(3):
                label = table.cell(row_idx, 0).text.strip()
                if data.incident_category and label in data.incident_category:
                    set_cell_text(table, row_idx, 1, "X")
                else:
                    set_cell_text(table, row_idx, 1, "")  # Clear if not selected
        
        # Table 3 (6x2): Object Details
        elif rows == 6 and cols == 2:
            # Row 0: Specifics of Object
            set_cell_text(table, 0, 1, data.object_details.specifics)
            set_cell_text(table, 1, 1, data.object_details.size)
            set_cell_text(table, 2, 1, data.object_details.shape)
            set_cell_text(table, 3, 1, data.object_details.weight)
            set_cell_text(table, 4, 1, data.object_details.distance_reach)
            set_cell_text(table, 5, 1, data.object_details.item_name)


def fill_slide_3(slide, data: RCAIncident):
    """Fill Slide 3: Incident Description."""
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        
        table = shape.table
        rows = len(table.rows)
        
        # Table 2 (4x1): Description fields
        if rows == 4:
            set_cell_text(table, 1, 0, data.injury_description)
            # Use manager's account for RCA PowerPoint (not associate's account from WMW-738)
            set_cell_text(table, 3, 0, data.manager_incident_account or data.how_incident_occurred)


def add_images_to_slide(slide, image_paths: list[Path], max_images: int = 4):
    """Add images to a slide in a grid layout, maintaining aspect ratio."""
    if not image_paths:
        return
    
    from PIL import Image
    
    # Grid positions for up to 4 images (2x2 grid)
    # (left, top, max_width, max_height)
    positions = [
        (Inches(0.5), Inches(1.2), Inches(5.5), Inches(2.8)),   # Top-left
        (Inches(6.5), Inches(1.2), Inches(5.5), Inches(2.8)),   # Top-right
        (Inches(0.5), Inches(4.2), Inches(5.5), Inches(2.8)),   # Bottom-left
        (Inches(6.5), Inches(4.2), Inches(5.5), Inches(2.8)),   # Bottom-right
    ]
    
    for i, img_path in enumerate(image_paths[:max_images]):
        if img_path.exists():
            left, top, max_width, max_height = positions[i]
            
            # Get image dimensions to calculate aspect ratio
            try:
                with Image.open(img_path) as img:
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height
                    
                    # Calculate dimensions that fit within max bounds while maintaining aspect ratio
                    max_w_inches = max_width.inches
                    max_h_inches = max_height.inches
                    
                    # Try fitting to width first
                    new_width = max_w_inches
                    new_height = new_width / aspect_ratio
                    
                    # If height exceeds max, fit to height instead
                    if new_height > max_h_inches:
                        new_height = max_h_inches
                        new_width = new_height * aspect_ratio
                    
                    # Center the image within the max bounds
                    actual_left = left + Inches((max_w_inches - new_width) / 2)
                    actual_top = top + Inches((max_h_inches - new_height) / 2)
                    
                    slide.shapes.add_picture(
                        str(img_path), 
                        actual_left, 
                        actual_top, 
                        Inches(new_width), 
                        Inches(new_height)
                    )
            except Exception as e:
                # Fallback: add with just width constraint
                slide.shapes.add_picture(str(img_path), left, top, width=max_width)


def fill_slide_6(slide, data: RCAIncident):
    """Fill Slide 6: Root Cause Analysis (5 Whys with 4Ms matrix)."""
    five_whys = data.five_whys
    
    for shape in slide.shapes:
        # Fill Problem Statement (TextBox 6)
        if shape.has_text_frame:
            if "TextBox 6" in shape.name or (hasattr(shape, 'text') and shape.text == ""):
                # Look for the problem statement text box
                if shape.name == "TextBox 6":
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.text = five_whys.problem_statement
                        for run in paragraph.runs:
                            run.font.color.rgb = BLACK
        
        # Fill 5 Whys Table (Table 7 - 5 rows x 7 cols)
        if shape.has_table:
            table = shape.table
            rows = len(table.rows)
            cols = len(table.columns)
            
            # Check if this is the 5 Whys table (5 rows x 7 cols)
            if rows == 5 and cols == 7:
                # Row 1: Material
                set_cell_text(table, 1, 1, five_whys.material.why1)
                set_cell_text(table, 1, 2, five_whys.material.why2)
                set_cell_text(table, 1, 3, five_whys.material.why3)
                set_cell_text(table, 1, 4, five_whys.material.why4)
                set_cell_text(table, 1, 5, five_whys.material.why5)
                set_cell_text(table, 1, 6, five_whys.material.root_cause)
                
                # Row 2: Machine
                set_cell_text(table, 2, 1, five_whys.machine.why1)
                set_cell_text(table, 2, 2, five_whys.machine.why2)
                set_cell_text(table, 2, 3, five_whys.machine.why3)
                set_cell_text(table, 2, 4, five_whys.machine.why4)
                set_cell_text(table, 2, 5, five_whys.machine.why5)
                set_cell_text(table, 2, 6, five_whys.machine.root_cause)
                
                # Row 3: Method
                set_cell_text(table, 3, 1, five_whys.method.why1)
                set_cell_text(table, 3, 2, five_whys.method.why2)
                set_cell_text(table, 3, 3, five_whys.method.why3)
                set_cell_text(table, 3, 4, five_whys.method.why4)
                set_cell_text(table, 3, 5, five_whys.method.why5)
                set_cell_text(table, 3, 6, five_whys.method.root_cause)
                
                # Row 4: huMan
                set_cell_text(table, 4, 1, five_whys.human.why1)
                set_cell_text(table, 4, 2, five_whys.human.why2)
                set_cell_text(table, 4, 3, five_whys.human.why3)
                set_cell_text(table, 4, 4, five_whys.human.why4)
                set_cell_text(table, 4, 5, five_whys.human.why5)
                set_cell_text(table, 4, 6, five_whys.human.root_cause)


def fill_slide_7(slide, data: RCAIncident):
    """Fill Slide 7: Countermeasures."""
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        
        table = shape.table
        rows = len(table.rows)
        cols = len(table.columns)
        
        # Table 15 (5x6): Countermeasures table
        if rows == 5 and cols == 6:
            for i, cm in enumerate(data.countermeasures[:4]):
                row_idx = i + 1  # Skip header row
                # Only fill row if there's actual content (root cause or countermeasure)
                has_content = cm.root_cause or cm.countermeasure
                set_cell_text(table, row_idx, 0, cm.root_cause if has_content else "")
                set_cell_text(table, row_idx, 1, cm.countermeasure if has_content else "")
                set_cell_text(table, row_idx, 2, cm.owner if has_content else "")
                # Only show due date and pyramid level if there's content
                set_cell_text(table, row_idx, 3, format_date(cm.due_date) if has_content and cm.due_date else "")
                set_cell_text(table, row_idx, 4, cm.cost if has_content else "")
                set_cell_text(table, row_idx, 5, str(cm.pyramid_level) if has_content and cm.pyramid_level else "")


def generate_rca_pptx(
    data: RCAIncident,
    template_path: Path,
    output_path: Path,
    reenactment_photos: list[Path] = None,
    cctv_screenshots: list[Path] = None,
) -> Path:
    """
    Generate a filled RCA PowerPoint from template.
    
    Args:
        data: The form data
        template_path: Path to RCA.pptx template
        output_path: Where to save the filled PPTX
        reenactment_photos: List of paths to reenactment photos
        cctv_screenshots: List of paths to CCTV screenshots
    
    Returns:
        Path to the generated PPTX file
    """
    # Copy template to output location first
    shutil.copy(template_path, output_path)
    
    # Open and modify
    prs = Presentation(output_path)
    slides = list(prs.slides)
    
    # Fill each slide
    if len(slides) >= 1:
        fill_slide_1(slides[0], data)
    
    if len(slides) >= 2:
        fill_slide_2(slides[1], data)
    
    if len(slides) >= 3:
        fill_slide_3(slides[2], data)
    
    if len(slides) >= 4 and reenactment_photos:
        add_images_to_slide(slides[3], reenactment_photos)
    
    if len(slides) >= 5 and cctv_screenshots:
        add_images_to_slide(slides[4], cctv_screenshots)
    
    if len(slides) >= 6:
        fill_slide_6(slides[5], data)
    
    if len(slides) >= 7:
        fill_slide_7(slides[6], data)
    
    # Save
    prs.save(output_path)
    return output_path
